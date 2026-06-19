from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/Dyslexia_Detection_System_Demo_Deck.pptx")

WIDE = (13.333, 7.5)
NAVY = RGBColor(14, 45, 84)
TEAL = RGBColor(25, 153, 145)
GOLD = RGBColor(215, 164, 72)
SLATE = RGBColor(73, 84, 99)
LIGHT = RGBColor(245, 247, 250)
MID = RGBColor(108, 117, 125)


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])


def add_slide_base(slide, title: str, num: int) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.22), Inches(1.35), Inches(0.12)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(11.9), Inches(0.65))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY

    num_box = slide.shapes.add_textbox(Inches(12.0), Inches(6.95), Inches(0.7), Inches(0.25))
    tfn = num_box.text_frame
    tfn.clear()
    pn = tfn.paragraphs[0]
    pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run()
    rn.text = f"{num}"
    rn.font.size = Pt(10)
    rn.font.color.rgb = MID


def add_text_slide(prs, slide, title: str, bullets, num: int, footer: str | None = None):
    add_slide_base(slide, title, num)
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.8), Inches(5.9))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(18 if len(item) < 80 else 16)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)
        p.bullet = True
    if footer:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(6.45), Inches(11.0), Inches(0.3))
        tf2 = box.text_frame
        p = tf2.paragraphs[0]
        r = p.add_run()
        r.text = footer
        r.font.size = Pt(11)
        r.font.italic = True
        r.font.color.rgb = MID


def add_two_column_slide(prs, slide, title: str, left_title: str, left_items, right_title: str, right_items, num: int):
    add_slide_base(slide, title, num)
    for x, heading, items in [
        (0.7, left_title, left_items),
        (6.8, right_title, right_items),
    ]:
        head = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(5.4), Inches(0.4))
        p = head.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.bold = True
        r.font.size = Pt(17)
        r.font.color.rgb = NAVY
        box = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(5.4), Inches(4.9))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = item
            p.level = 0
            p.font.name = "Aptos"
            p.font.size = Pt(15)
            p.font.color.rgb = SLATE
            p.space_after = Pt(5)
            p.bullet = True


def add_table_slide(prs, slide, title: str, rows, num: int, col_widths=None):
    add_slide_base(slide, title, num)
    left = Inches(0.65)
    top = Inches(1.35)
    width = Inches(12.0)
    height = Inches(4.7)
    table = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(13 if r == 0 else 12)
                p.font.color.rgb = NAVY if r == 0 else SLATE
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = LIGHT if r == 0 else RGBColor(255, 255, 255)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    note = slide.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(11.5), Inches(0.4))
    p = note.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Prepared from the current repository documentation and project flow."
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = MID


def add_cover(prs, slide, num: int):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.25))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()

    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.7), Inches(0.8), Inches(2.0), Inches(2.0))
    orb.fill.solid()
    orb.fill.fore_color.rgb = TEAL
    orb.fill.transparency = 0.2
    orb.line.fill.background()

    orb2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.2), Inches(4.8), Inches(2.8), Inches(2.8))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = GOLD
    orb2.fill.transparency = 0.35
    orb2.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(9.5), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Dyslexia Detection System"
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(9.6), Inches(1.1))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "28-slide demonstration deck for screening, therapy, biomarkers, reports, and research planning"
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(232, 238, 244)

    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.8), Inches(6.5), Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(255, 255, 255)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Built from the project's architecture, model catalog, research draft, and dashboard documentation."
    r.font.size = Pt(16)
    r.font.color.rgb = SLATE

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(7.5), Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Demo presentation for local use"
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(220, 228, 235)


prs = Presentation()
set_wide_layout(prs)

slides = []

slides.append(("cover", None))
slides.append(("text", ("Agenda", [
    "Problem, scope, and system overview",
    "Data inputs, preprocessing, and model families",
    "Screening, therapy, biomarkers, and explanation modules",
    "Dashboards, records, evaluation, and research directions",
])))
slides.append(("text", ("Problem Statement", [
    "Dyslexia signals often appear across handwriting, speech, reading behavior, and eye movement rather than one isolated feature.",
    "Single-modality screening can miss important cues and can be difficult to explain to teachers or parents.",
    "The project is designed to support local-first, interpretable, multilingual educational assistance.",
])))
slides.append(("text", ("Why Multimodal Screening", [
    "Handwriting captures shape and spacing cues.",
    "Speech captures fluency, pauses, and pronunciation quality.",
    "Text captures reading and spelling patterns.",
    "Behavior and gaze capture attention, hesitation, and visual search strategy.",
])))
slides.append(("text", ("Project Scope", [
    "Target users: students, teachers, parents, clinicians, and researchers.",
    "Supported tasks: screening, speech therapy, visual focus testing, biomarker analysis, and report generation.",
    "Supported languages and views: English, Bengali, and extendable multilingual workflows.",
    "Deployment focus: local browser use, local persistence, and research-ready export.",
])))
slides.append(("text", ("System Overview", [
    "Input ingestion from manifest CSV, uploaded files, and live dashboard activities.",
    "Preprocessing before branch-specific feature extraction.",
    "Modality encoders for handwriting, audio, text, and behavior signals.",
    "Fusion and prediction followed by explanation, intervention, and reporting.",
])))
slides.append(("text", ("Data Sources and Manifest Contract", [
    "Central CSV manifest keeps samples portable and reproducible.",
    "Each sample can reference handwriting images, audio files, text samples, and behavior values.",
    "Optional eye-tracking and biomarker tables extend the same workflow.",
    "The manifest acts as the single source of truth for records and experiments.",
])))
slides.append(("text", ("Preprocessing Pipeline", [
    "Handwriting: grayscale loading, contrast normalization, and resizing.",
    "Audio: spectral feature extraction and stable frame preparation.",
    "Text: tokenization, padding, and language-specific vocabulary handling.",
    "Numeric features: cleaning, scaling, and shape validation.",
])))
slides.append(("table", ("Model Catalog Overview", [
    ["Component", "Purpose", "Main Advantage"],
    ["HandwritingEncoder", "Image feature learning", "Fast local inference"],
    ["AudioEncoder", "Audio feature learning", "Compact temporal cues"],
    ["TransformerTextEncoder", "Text context modeling", "Stronger sequence learning"],
    ["AttentionFusionClassifier", "Modal weighting", "Interpretability"],
    ["InterventionPolicy", "Action selection", "Editable support logic"],
])))
slides.append(("text", ("Default Multimodal Screening Model", [
    "The default model combines handwriting, audio, text, and behavior embeddings.",
    "It is designed to be practical, balanced, and easy to deploy locally.",
    "The model outputs risk or severity estimates for downstream interpretation.",
])))
slides.append(("text", ("Attention Fusion", [
    "Attention fusion learns which modality contributes most for a given learner.",
    "This supports explanations such as handwriting-heavy, speech-heavy, or behavior-heavy cases.",
    "The output is useful for both model debugging and educator-facing summaries.",
])))
slides.append(("text", ("Transformer Multimodal Variant", [
    "Transformer-based branches improve contextual modeling for text and longer sequences.",
    "This variant is useful when sequence dependencies matter more than simple recurrence.",
    "It is heavier than the baseline but can improve representation quality in research settings.",
])))
slides.append(("text", ("Vision Transformer Branch", [
    "Vision transformer variants split handwriting images into patches.",
    "Patch attention helps capture layout, spacing, and structural consistency patterns.",
    "This branch is suitable for experiments that need stronger image reasoning than CNN baselines.",
])))
slides.append(("text", ("Speech and Therapy Workflow", [
    "Speech therapy mode listens to a reading sample, transcribes it, and scores reading behavior.",
    "The workflow can surface current prompt, recognized response, spelling score, and overall status.",
    "It is intended for guided practice, not only for assessment.",
])))
slides.append(("text", ("Text and Spelling Branch", [
    "Character-level text modeling supports short prompts, spelling practice, and reading-response analysis.",
    "The branch can be used alongside audio to compare spoken and written performance.",
    "It also helps keep the system multilingual and lightweight.",
])))
slides.append(("text", ("Behavior and Eye-Tracking Signals", [
    "Behavior metrics include hesitation, repetition, omission, reading speed, and timing values.",
    "Eye-tracking metrics help measure fixation, gaze control, and visual attention consistency.",
    "These features strengthen the low-resource signal set when media inputs are limited.",
])))
slides.append(("text", ("Biomarker Discovery", [
    "The biomarkers page accepts structured CSV tables for analysis.",
    "The system can detect likely label columns and highlight strong numeric signals.",
    "This is useful for simple feature ranking and plain-language summaries.",
])))
slides.append(("text", ("Screening Output and Severity Interpretation", [
    "The screening result summarizes language, predicted severity, confidence, and evidence by modality.",
    "The output is meant to read like a short educational summary, not just a raw class label.",
    "Different views can be tuned for classroom, home, and learner-friendly messaging.",
])))
slides.append(("text", ("Explainability Outputs", [
    "Heatmaps and attention scores help users see why the system reached a result.",
    "Model-level explanations can be combined with human-readable summaries.",
    "Explainability is central to teacher and parent trust in the workflow.",
])))
slides.append(("text", ("Speech Therapy Module", [
    "Shows prompt, response, auto score, passing threshold, and overall status.",
    "Supports round-based guided reading with automatic progression.",
    "Useful for repeated practice and progress tracking across sessions.",
])))
slides.append(("text", ("Visual Focus Test", [
    "The visual focus test supports letters, digits, and mixed-symbol modes.",
    "Each round asks the learner to tap a target symbol from an answer grid.",
    "The module is designed to measure attention, speed, accuracy, and consistency.",
])))
slides.append(("text", ("Test Lab and Report", [
    "The test and report area collects learner details before generating the final report.",
    "It keeps the output ready for PDF generation and session documentation.",
    "This page is the bridge between live testing and formal reporting.",
])))
slides.append(("text", ("Records and Local History", [
    "Saved sessions build a local history of tests and reports.",
    "The dashboard can show total records, filtered results, latest entry, and last saved time.",
    "This helps the system behave like a continuing local toolkit instead of a one-off demo.",
])))
slides.append(("table", ("Evaluation Metrics", [
    ["Metric", "What It Measures", "Why It Matters"],
    ["Accuracy / F1", "Prediction quality", "Basic model reliability"],
    ["Confidence", "How sure the system is", "Useful for risk triage"],
    ["Stability", "Consistency across runs", "Important for trust"],
    ["Sensitivity", "Detection of difficult cases", "Important for early support"],
    ["Latency", "Runtime on local devices", "Important for practical use"],
])))
slides.append(("two", ("Deployment Surfaces", "Streamlit dashboard", [
    "Fast research exploration and experimentation",
    "Useful for single-user local workflows",
    "Good for quick checks and model validation",
], "Web dashboard", [
    "Browser-based local UI",
    "Supports live interaction and record keeping",
    "Best for demonstration and sharing",
])))
slides.append(("text", ("Research Opportunities", [
    "Collect larger multilingual datasets.",
    "Add stronger fusion methods and self-supervised pretraining.",
    "Validate intervention suggestions in classroom-like settings.",
    "Extend explanation quality for teachers, parents, and learners.",
])))
slides.append(("text", ("Limitations and Ethics", [
    "The system is a support tool, not a diagnosis replacement.",
    "Data quality, consent, and privacy remain essential.",
    "Model outputs should be reviewed by trained professionals in real educational settings.",
])))
slides.append(("text", ("Conclusion", [
    "This project unifies screening, therapy, biomarkers, reporting, and records in one local-first platform.",
    "The next step is to expand evaluation, strengthen data collection, and validate educational usefulness.",
    "The deck is ready for demonstration, review, and future publication work.",
])))

slide_index = 1
for kind, payload in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if kind == "cover":
        add_cover(prs, slide, slide_index)
    elif kind == "text":
        title, bullets = payload
        footer = None
        if title == "Conclusion":
            footer = "Thank you"
        add_text_slide(prs, slide, title, bullets, slide_index, footer=footer)
    elif kind == "two":
        title, left_title, left_items, right_title, right_items = payload
        add_two_column_slide(prs, slide, title, left_title, left_items, right_title, right_items, slide_index)
    elif kind == "table":
        title, rows = payload
        widths = None
        if title == "Evaluation Metrics":
            widths = [2.0, 4.2, 5.8]
        elif title == "Model Catalog Overview":
            widths = [3.4, 4.6, 3.8]
        else:
            widths = [2.4, 4.2, 5.4]
        add_table_slide(prs, slide, title, rows, slide_index, col_widths=widths)
    slide_index += 1

OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT_PATH))
print(f"Saved {OUT_PATH} with {len(prs.slides)} slides")
